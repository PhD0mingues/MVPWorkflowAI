#!/usr/bin/env python3
"""
Monta a proposta comercial final como .pptx, no estilo visual da Advisia.

Uso:
    python3 montar_proposta_pptx.py dados_proposta.json saida.pptx [--estilo estilo.json]

Onde dados_proposta.json tem o formato:
{
  "cliente": "Nome do Cliente",
  "data": "08/07/2026",
  "secoes_criativas": {
    "abertura": "texto...",
    "diagnostico": "texto...",
    "proposta_de_valor": "texto..."
  },
  "campos_deterministicos": {
    "escopo": "texto vindo da tabela, ou [PENDENTE — preencher com tabela oficial]",
    "preco": "[PENDENTE — preencher com tabela oficial]",
    "prazo": "[PENDENTE — preencher com tabela oficial]",
    "condicoes_pagamento": "[PENDENTE — preencher com tabela oficial]"
  },
  "templates_usados": ["proposta-123: selecionada por correlação de escopo de automação"]
}

E estilo.json (opcional) tem o mesmo formato usado pela skill correlacao-projetos-ppt,
gerado por scripts/extrair_estilo_pptx.py daquela skill (disponível em
"$CLAUDE_PLUGIN_ROOT/skills/correlacao-projetos-ppt/scripts/extrair_estilo_pptx.py"
dentro deste mesmo plugin) a partir do PPT-modelo real de propostas:
{
  "fonte": "Georgia",
  "cor_barra_titulo": "1F1F4A",
  "cor_barra_cluster": "C99A2E",
  "cor_texto_titulo": "1F1F4A",
  "cor_texto_corpo": "333333",
  "cor_texto_rodape": "C99A2E",
  "cor_texto_barra": "FFFFFF"
}

Por que --estilo em vez de cor fixa: o objetivo é sempre usar as cores REAIS
do PPT-modelo de propostas, lidas diretamente do arquivo. Isso só é possível
quando você tem os bytes do arquivo (upload local, ou pasta local/OneDrive) -
rode extrair_estilo_pptx.py nele primeiro. Quando a única fonte disponível é
texto vindo do conector Microsoft 365 (que não expõe bytes/tema), não tem
como ler a cor real; nesse caso o script cai nos valores abaixo (última
cópia conhecida das cores reais da Advisia), avisando que podem estar
desatualizadas.

Não adivinhe conteúdo aqui - este script só aplica o estilo visual (cores,
fontes, posições). O texto das seções criativas e os campos determinísticos
vêm sempre dos Passos 5 e 6 do SKILL.md de gerar-proposta.
"""
import argparse
import json
import sys

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor

# Última cópia conhecida de cores REAIS do template Advisia (mesma paleta
# usada pela skill correlacao-projetos-ppt). Usada só como fallback quando
# --estilo não é passado.
ESTILO_PADRAO = {
    "fonte": "Georgia",
    "cor_barra_titulo": "1F1F4A",
    "cor_barra_cluster": "C99A2E",
    "cor_texto_titulo": "1F1F4A",
    "cor_texto_corpo": "333333",
    "cor_texto_rodape": "C99A2E",
    "cor_texto_barra": "FFFFFF",
    "cor_alerta": "B00000",
}

SIDEBAR_W = Emu(2926080)
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
CONTENT_X = Emu(3291840)
CONTENT_W = Emu(8229600)

SECAO_TITULOS = {
    "abertura": "Abertura",
    "diagnostico": "Diagnóstico",
    "proposta_de_valor": "Proposta de Valor",
}

CAMPO_TITULOS = {
    "escopo": "Escopo",
    "preco": "Investimento",
    "prazo": "Prazo",
    "condicoes_pagamento": "Condições de Pagamento",
}

PENDENTE_MARCADOR = "[PENDENTE"


def carregar_estilo(caminho):
    if not caminho:
        return dict(ESTILO_PADRAO)
    with open(caminho, encoding="utf-8") as f:
        lido = json.load(f)
    estilo = dict(ESTILO_PADRAO)
    for chave, valor in lido.items():
        if valor:
            estilo[chave] = valor
    return estilo


def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def add_textbox(slide, x, y, w, h, text, size, bold, color, font, align="l"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {"l": 1, "ctr": 2}.get(align, 1)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    return box


def add_sidebar(slide, color):
    rect = slide.shapes.add_shape(1, Emu(0), Emu(0), SIDEBAR_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def build_title_slide(prs, estilo, cliente, data):
    font = estilo["fonte"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_sidebar(slide, hex_to_rgb(estilo["cor_barra_titulo"]))
    add_textbox(slide, CONTENT_X, Emu(2377440), CONTENT_W, Emu(1097280),
                f"Proposta Comercial — {cliente}", 36, True,
                hex_to_rgb(estilo["cor_texto_titulo"]), font)
    add_textbox(slide, CONTENT_X, Emu(3291840), CONTENT_W, Emu(548640),
                f"Data: {data}", 16, False, hex_to_rgb(estilo["cor_texto_corpo"]), font)
    return slide


def build_section_slide(prs, estilo, titulo, texto):
    font = estilo["fonte"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_sidebar(slide, hex_to_rgb(estilo["cor_barra_cluster"]))
    add_textbox(slide, CONTENT_X, Emu(640080), CONTENT_W, Emu(914400),
                titulo, 28, True, hex_to_rgb(estilo["cor_texto_titulo"]), font)
    add_textbox(slide, CONTENT_X, Emu(1737360), CONTENT_W, Emu(4389120),
                texto, 16, False, hex_to_rgb(estilo["cor_texto_corpo"]), font)
    return slide


def build_condicoes_slide(prs, estilo, campos):
    font = estilo["fonte"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_sidebar(slide, hex_to_rgb(estilo["cor_barra_cluster"]))
    add_textbox(slide, CONTENT_X, Emu(640080), CONTENT_W, Emu(914400),
                "Escopo, Investimento e Condições", 28, True,
                hex_to_rgb(estilo["cor_texto_titulo"]), font)

    y = 1737360
    pendentes = []
    for chave, titulo in CAMPO_TITULOS.items():
        valor = campos.get(chave, f"{PENDENTE_MARCADOR} — não fornecido]")
        eh_pendente = PENDENTE_MARCADOR in str(valor)
        if eh_pendente:
            pendentes.append(titulo)
        cor = hex_to_rgb(estilo["cor_alerta"]) if eh_pendente else hex_to_rgb(estilo["cor_texto_corpo"])
        add_textbox(slide, CONTENT_X, Emu(y), CONTENT_W, Emu(548640),
                    f"{titulo}: {valor}", 16, eh_pendente, cor, font)
        y += 548640

    if pendentes:
        add_textbox(slide, CONTENT_X, Emu(y + 100000), CONTENT_W, Emu(731520),
                    "ATENÇÃO: campo(s) pendente(s) de tabela oficial — não enviar ao "
                    "cliente antes de resolver.", 14, True,
                    hex_to_rgb(estilo["cor_alerta"]), font)
    return slide


def build_notas_slide(prs, estilo, templates_usados):
    if not templates_usados:
        return None
    font = estilo["fonte"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_sidebar(slide, hex_to_rgb(estilo["cor_barra_cluster"]))
    add_textbox(slide, CONTENT_X, Emu(640080), CONTENT_W, Emu(914400),
                "Nota interna — templates de referência usados", 24, True,
                hex_to_rgb(estilo["cor_texto_titulo"]), font)
    texto = "\n".join(f"• {t}" for t in templates_usados)
    add_textbox(slide, CONTENT_X, Emu(1737360), CONTENT_W, Emu(4389120),
                texto, 14, False, hex_to_rgb(estilo["cor_texto_corpo"]), font)
    return slide


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("dados_json")
    parser.add_argument("saida_pptx")
    parser.add_argument(
        "--estilo",
        help="JSON de estilo gerado por extrair_estilo_pptx.py a partir do PPT-modelo real de propostas. "
        "Se omitido, usa a última cópia conhecida das cores reais da Advisia.",
    )
    args = parser.parse_args()

    with open(args.dados_json, encoding="utf-8") as f:
        dados = json.load(f)

    estilo = carregar_estilo(args.estilo)
    if not args.estilo:
        print(
            "[aviso] Rodando sem --estilo - usando última cópia conhecida das cores "
            "reais da Advisia, não uma leitura ao vivo do PPT-modelo.",
            file=sys.stderr,
        )

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs, estilo, dados.get("cliente", "[cliente não informado]"), dados.get("data", ""))

    secoes = dados.get("secoes_criativas", {})
    for chave, titulo in SECAO_TITULOS.items():
        if chave in secoes:
            build_section_slide(prs, estilo, titulo, secoes[chave])

    build_condicoes_slide(prs, estilo, dados.get("campos_deterministicos", {}))
    build_notas_slide(prs, estilo, dados.get("templates_usados", []))

    prs.save(args.saida_pptx)
    print(f"Proposta salva em {args.saida_pptx}")

    campos = dados.get("campos_deterministicos", {})
    pendentes = [k for k, v in campos.items() if PENDENTE_MARCADOR in str(v)]
    if pendentes:
        print(
            f"ATENÇÃO: campos determinísticos pendentes: {', '.join(pendentes)}. "
            "Não considere o fluxo concluído sem resolver isso.",
            file=sys.stderr,
        )


if __name__ == "__main__":
    main()
