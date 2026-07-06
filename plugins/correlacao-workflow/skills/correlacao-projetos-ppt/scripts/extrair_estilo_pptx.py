#!/usr/bin/env python3
"""
Le um .pptx de referencia (bytes reais, nao texto extraido) e extrai o
estilo visual usado por ele: cores de barra lateral, cores de texto e
fonte. Feito sob medida para o padrao de template da Advisia (barra
lateral colorida a esquerda + textboxes soltos), mas funciona de forma
heuristica em qualquer .pptx parecido.

Uso:
    python3 extrair_estilo_pptx.py referencia.pptx > estilo.json

So roda se voce tiver os bytes reais do arquivo (upload local, ou pasta
local/OneDrive) - nao funciona a partir de texto vindo do conector M365,
porque nao ha bytes pra abrir com python-pptx nesse caso.
"""
import json
import sys

from pptx import Presentation
from pptx.util import Emu


def rgb_hex(color):
    try:
        return str(color.rgb)
    except Exception:
        return None


def find_sidebar_color(slide):
    best = None
    for shape in slide.shapes:
        if shape.left == 0 and shape.top == 0:
            try:
                if shape.fill.type == 1:
                    color = rgb_hex(shape.fill.fore_color)
                    if color:
                        if best is None or shape.height > best[1]:
                            best = (color, shape.height)
            except Exception:
                continue
    return best[0] if best else None


def collect_runs(slide):
    runs = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                size = run.font.size.pt if run.font.size else None
                color = None
                try:
                    color = rgb_hex(run.font.color)
                except Exception:
                    pass
                runs.append(
                    {
                        "text": run.text,
                        "size": size,
                        "bold": bool(run.font.bold),
                        "color": color,
                        "font": run.font.name,
                    }
                )
    return runs


def main():
    if len(sys.argv) != 2:
        print("Uso: python3 extrair_estilo_pptx.py referencia.pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(sys.argv[1])
    slides = list(prs.slides)
    if not slides:
        print("PPT sem slides.", file=sys.stderr)
        sys.exit(1)

    estilo = {
        "fonte": None,
        "cor_barra_titulo": None,
        "cor_barra_cluster": None,
        "cor_texto_titulo": None,
        "cor_texto_corpo": None,
        "cor_texto_rodape": None,
        "cor_texto_barra": None,
    }

    slide1 = slides[0]
    estilo["cor_barra_titulo"] = find_sidebar_color(slide1)
    runs1 = collect_runs(slide1)
    if runs1:
        runs_by_size = sorted([r for r in runs1 if r["size"]], key=lambda r: -r["size"])
        if runs_by_size:
            estilo["cor_texto_titulo"] = runs_by_size[0]["color"]
            estilo["fonte"] = runs_by_size[0]["font"]
        corpo = [r for r in runs_by_size if not r["bold"]]
        if corpo:
            estilo["cor_texto_corpo"] = corpo[0]["color"]
        rodape_candidatos = sorted(
            [r for r in runs1 if r["bold"] and r["size"]], key=lambda r: r["size"]
        )
        if rodape_candidatos:
            estilo["cor_texto_rodape"] = rodape_candidatos[0]["color"]

    if len(slides) > 1:
        slide2 = slides[1]
        estilo["cor_barra_cluster"] = find_sidebar_color(slide2)
        barra_runs = []
        for shape in slide2.shapes:
            if shape.has_text_frame and shape.left is not None and shape.left < Emu(2926080):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            try:
                                barra_runs.append(rgb_hex(run.font.color))
                            except Exception:
                                pass
        if barra_runs:
            estilo["cor_texto_barra"] = barra_runs[0]

    print(json.dumps(estilo, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
