#!/usr/bin/env python3
"""
Extrai todo o texto (títulos, corpo, notas do apresentador) de cada .pptx
encontrado dentro de uma pasta raiz, recursivamente, e imprime um JSON
com uma entrada por arquivo encontrado.

Uso:
    python3 extract_pptx_text.py <pasta_raiz> [--max-depth N]

Saída (stdout), lista de objetos:
    [
      {"project": "1", "file": "Workflows AI/1/1.pptx", "text": "Cigarro"},
      ...
    ]

"project" é o nome da subpasta imediatamente acima do arquivo (ou o nome
do arquivo sem extensão, se o .pptx estiver direto na pasta raiz).
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation


def extract_text_from_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    chunks = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_chunks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_chunks.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        t = cell.text.strip()
                        if t:
                            slide_chunks.append(t)
        # notas do apresentador
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_chunks.append(f"[notas] {notes}")
        if slide_chunks:
            chunks.append(f"--- slide {slide_idx} ---\n" + "\n".join(slide_chunks))
    return "\n".join(chunks)


def find_pptx_files(root: Path):
    return sorted(root.rglob("*.pptx"))


def project_name_for(path: Path, root: Path) -> str:
    try:
        rel = path.relative_to(root)
    except ValueError:
        rel = path
    parts = rel.parts
    if len(parts) > 1:
        return parts[0]
    return path.stem


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("root", help="Pasta raiz contendo as pastas de projeto")
    args = parser.parse_args()

    root = Path(args.root).expanduser().resolve()
    if not root.exists():
        print(f"Pasta não encontrada: {root}", file=sys.stderr)
        sys.exit(1)

    files = find_pptx_files(root)
    if not files:
        print(f"Nenhum .pptx encontrado em {root}", file=sys.stderr)
        sys.exit(1)

    results = []
    for f in files:
        # ignora pptx dentro de uma subpasta "outputs" (resultado de execuções anteriores)
        if "outputs" in f.relative_to(root).parts[:-1]:
            continue
        try:
            text = extract_text_from_pptx(f)
        except Exception as e:  # noqa: BLE001
            text = f"[ERRO ao ler arquivo: {e}]"
        results.append(
            {
                "project": project_name_for(f, root),
                "file": str(f.relative_to(root)),
                "text": text,
            }
        )

    print(json.dumps(results, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
