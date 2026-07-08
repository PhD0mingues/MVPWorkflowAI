#!/usr/bin/env python3
"""
Monta o PPT de resultado da skill correlacao-projetos-ppt.

Uso:
    python3 montar_correlacoes_pptx.py dados.json saida.pptx [--estilo estilo.json]

Onde dados.json tem o formato:
{
  "titulo": "Correlações entre Projetos",
  "subtitulo": "Leitura orientada ao mercado de design e publicidade de eletrodomésticos",
  "rodape_data": "Workflows AI — 03/07/2026",
  "clusters": [
    {
      "nome": "O aroma como atributo de design",
      "itens": ["6 - Cheiro", "3 - Geladeira", "4 - Forno", "1 - Cigarro"],
      "explicacao": "texto da correlação..."
    }
  ]
}

E estilo.json (opcional, gerado por scripts/extrair_estilo_pptx.py a partir
de um .pptx de referência real) tem o formato:
{
  "fonte": "Georgia",
  "cor_barra_titulo": "1F1F4A",
  "cor_barra_cluster": "C99A2E",
  "cor_texto_titulo": "1F1F4A",
  "cor_texto_corpo": "333333",
  "cor_texto_rodape": "C99A2E",
  "cor_texto_barra": "FFFFFF"
}

Por que existe --estilo em vez de cor fixa: o objetivo e sempre usar as
cores REAIS do template de referencia, lidas diretamente do arquivo. Isso so
e possivel quando voce tem os bytes do arquivo (upload local, ou pasta
local/OneDrive) - rode extrair_estilo_pptx.py nele primeiro. Quando a unica
fonte disponivel e texto vindo do conector Microsoft 365 (que nao expoe
bytes/tema), nao tem como ler a cor real; nesse caso o script cai nos valores
abaixo, que sao a ultima copia de cores REAIS conhecidas do template da
Advisia (extraidas de bytes reais em execucao anterior) - nao sao um "chute"
generico, mas tambem nao sao garantidas de bater com uma eventual atualizacao
visual do template. Sempre prefira rodar com --estilo quando tiver o arquivo.

Nao adivinhe conteudo aqui - este script so aplica o estilo visual (cores,
fontes, posicoes). Os clusters, nomes e explicacoes vem do raciocinio feito
nos Passos 3 e 4 do SKILL.md.
"""
import argparse
import json
import sys

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor

# Ultima copia conhecida de cores REAIS do template Advisia, extraidas de
# bytes reais do arquivo correlacoes_eletrodomesticos.pptx (nao inventadas).
# Usada so como fallback quando --estilo nao e passado.
ESTILO_PADRAO = {
    "fonte": "Georgia",
    "cor_barra_titulo": "1F1F4A",
    "cor_barra_cluster": "C99A2E",
    "cor_texto_titulo": "1F1F4A",
    "cor_texto_corpo": "333333",
    "cor_texto_rodape": "C99A2E",
    "cor_texto_barra": "FFFFFF",
}

SIDEBAR_W = Emu(2926080)
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
CONTENT_X = Emu(3291840)
CONTENT_W = Emu(8229600)


def carregar_estilo(caminho):
    if not caminho:
        return dict(ESTILO_PADRAO)
    with open(caminho, encoding="utf-8") as f:
        lido = json.load(f)
    estilo = dict(ESTILO_PADRAO)
    for chave, valor in lido.items():
        if valor:
            estilo[chave] = valor
    return estilo


def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def add_textbox(slide, x, y, w, h, text, size, bold, color, font, align="l"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {"l": 1, "ctr": 2}.get(align, 1)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    return box


def add_sidebar(slide, color):
    rect = slide.shapes.add_shape(1, Emu(0), Emu(0), SIDEBAR_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def build_title_slide(prs, estilo, titulo, subtitulo, rodape_data):
    font = estilo["fonte"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_sidebar(slide, hex_to_rgb(estilo["cor_barra_titulo"]))
    add_textbox(slide, CONTENT_X, Emu(2377440), CONTENT_W, Emu(1097280),
                titulo, 40, True, hex_to_rgb(estilo["cor_texto_titulo"]), font)
    add_textbox(slide, CONTENT_X, Emu(3291840), CONTENT_W, Emu(731520),
                subtitulo, 18, False, hex_to_rgb(estilo["cor_texto_corpo"]), font)
    add_textbox(slide, CONTENT_X, Emu(3840480), CONTENT_W, Emu(548640),
                rodape_data, 14, True, hex_to_rgb(estilo["cor_texto_rodape"]), font)
    return slide


def build_cluster_slide(prs, estilo, nome, itens, explicacao):
    font = estilo["fonte"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_sidebar(slide, hex_to_rgb(estilo["cor_barra_cluster"]))

    y = 640080
    for item in itens:
        add_textbox(slide, Emu(274320), Emu(y), Emu(2377440), Emu(457200),
                    item, 16, True, hex_to_rgb(estilo["cor_texto_barra"]), font)
        y += 457200

    add_textbox(slide, CONTENT_X, Emu(640080), CONTENT_W, Emu(914400),
                nome, 28, True, hex_to_rgb(estilo["cor_texto_titulo"]), font)
    add_textbox(slide, CONTENT_X, Emu(1737360), CONTENT_W, Emu(4389120),
                explicacao, 16, False, hex_to_rgb(estilo["cor_texto_corpo"]), font)
    return slide


def build_fontes_slide(prs, estilo, fontes):
    font = estilo["fonte"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_sidebar(slide, hex_to_rgb(estilo["cor_barra_cluster"]))

    add_textbox(slide, CONTENT_X, Emu(640080), CONTENT_W, Emu(731520),
                "Fontes utilizadas nesta análise (MVP)", 28, True,
                hex_to_rgb(estilo["cor_texto_titulo"]), font)

    y = Emu(1554480)
    line_h = Emu(365760)
    for fonte in fontes:
        add_textbox(slide, CONTENT_X, y, CONTENT_W, line_h,
                    f"- {fonte}", 14, False, hex_to_rgb(estilo["cor_texto_corpo"]), font)
        y = Emu(int(y) + int(line_h))
    return slide


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("dados_json")
    parser.add_argument("saida_pptx")
    parser.add_argument(
        "--estilo",
        help="JSON de estilo gerado por extrair_estilo_pptx.py a partir de um .pptx real. "
        "Se omitido, usa a ultima copia conhecida das cores reais da Advisia.",
    )
    args = parser.parse_args()

    with open(args.dados_json, encoding="utf-8") as f:
        dados = json.load(f)

    estilo = carregar_estilo(args.estilo)
    if not args.estilo:
        print(
            "[aviso] Rodando sem --estilo - usando ultima copia conhecida das cores "
            "reais da Advisia, nao uma leitura ao vivo do template.",
            file=sys.stderr,
        )

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(
        prs,
        estilo,
        dados.get("titulo", "Correlações entre Projetos"),
        dados.get("subtitulo", ""),
        dados.get("rodape_data", ""),
    )

    for cluster in dados.get("clusters", []):
        build_cluster_slide(
            prs,
            estilo,
            cluster["nome"],
            cluster.get("itens", []),
            cluster.get("explicacao", ""),
        )

    fontes = dados.get("fontes", [])
    if fontes:
        build_fontes_slide(prs, estilo, fontes)
    else:
        print(
            "[aviso] Nenhuma fonte informada em 'fontes' - o slide de "
            "transparencia (MVP) nao sera gerado.",
            file=sys.stderr,
        )

    prs.save(args.saida_pptx)
    print(f"PPT salvo em {args.saida_pptx}")


if __name__ == "__main__":
    main()
